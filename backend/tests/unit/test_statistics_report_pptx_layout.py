from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from app.core.statistics_report import build_pptx_report
from app.schemas.member_statistics import MembersStatistics

pytestmark = pytest.mark.unit


def _empty_stats() -> MembersStatistics:
    return MembersStatistics(
        by_status=[],
        by_join_year=[],
        growth_by_rotary_year=[],
        by_nationality=[],
        age_distribution=[],
        tenure_distribution=[],
        by_gender=[],
        total_members=0,
        honorary_members=0,
        new_members_this_rotary_year=0,
        countries_represented=0,
        women_count=0,
        men_count=0,
        average_age=None,
        average_tenure_as_rotarian=None,
    )


def test_title_heading_does_not_overlap_the_logo():
    pptx_bytes = build_pptx_report(_empty_stats())
    prs = Presentation(BytesIO(pptx_bytes))
    slide = prs.slides[0]

    picture = next(shape for shape in slide.shapes if shape.shape_type == 13)  # PICTURE
    heading_box = next(
        shape
        for shape in slide.shapes
        if shape.has_text_frame and "Members Statistics Report" in shape.text_frame.text
    )

    logo_right_edge = picture.left + picture.width
    assert heading_box.left >= logo_right_edge, (
        "Heading textbox starts before the logo's right edge — they overlap"
    )


def test_simplified_report_has_a_single_slide():
    pptx_bytes = build_pptx_report(_empty_stats(), report_type="simplified")
    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) == 1


def test_integral_report_adds_one_detail_slide_per_chart_dataset():
    # 6 chart datasets (join year, growth, nationality, tenure, gender, age)
    # => 1 main slide + 6 detail slides, per _detail_tables in statistics_report.py.
    pptx_bytes = build_pptx_report(_empty_stats(), report_type="integral")
    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) == 7


def test_template_slide_size_is_respected_and_content_still_fits():
    # A non-default (non-13.333in-wide) template should have its own slide
    # size preserved, with the from-scratch content scaled to it rather than
    # overflowing past its right/bottom edge.
    template_prs = Presentation()
    template_prs.slide_width = Inches(10)
    template_prs.slide_height = Inches(5.625)
    template_buf = BytesIO()
    template_prs.save(template_buf)
    template_buf.seek(0)

    template_path = Path("/tmp/test-ppt-template-unit.pptx")
    template_path.write_bytes(template_buf.getvalue())
    try:
        pptx_bytes = build_pptx_report(_empty_stats(), template_path=template_path)
        prs = Presentation(BytesIO(pptx_bytes))
        assert prs.slide_width == Inches(10)

        for shape in prs.slides[0].shapes:
            assert shape.left + shape.width <= prs.slide_width
            assert shape.top + shape.height <= prs.slide_height
    finally:
        template_path.unlink()


def _write_default_template(path: Path) -> None:
    # python-pptx's own default Presentation() ships 11 standard layouts,
    # including "Title Only" (a title placeholder, no body/object) — the
    # ideal candidate for Story 8.33's fix, and "Blank" (no placeholders at
    # all) — the layout the pre-8.33 code always picked, which is exactly
    # what silently discarded the template's own branding.
    Presentation().save(str(path))


class TestStory833TemplateBrandingFix:
    """Story 8.33 — an uploaded template's own title placeholder (and
    whatever master-level background/branding it inherits) must actually be
    used, not overridden by the app's own hardcoded blue/tone colours."""

    @pytest.fixture
    def template_path(self, tmp_path):
        path = tmp_path / "template.pptx"
        _write_default_template(path)
        return path

    def test_uses_the_layout_s_title_placeholder_instead_of_a_freehand_textbox(
        self, template_path
    ):
        pptx_bytes = build_pptx_report(_empty_stats(), template_path=template_path)
        prs = Presentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]

        assert slide.slide_layout.name == "Title Only"
        title_shapes = [shape for shape in slide.shapes if shape.is_placeholder]
        assert len(title_shapes) == 1
        assert "Members Statistics Report" in title_shapes[0].text_frame.text

    def test_does_not_force_the_app_s_brand_colour_onto_the_template_s_title(
        self, template_path
    ):
        pptx_bytes = build_pptx_report(_empty_stats(), template_path=template_path)
        prs = Presentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]

        title_placeholder = next(shape for shape in slide.shapes if shape.is_placeholder)
        # No explicit colour set on the paragraph — it's left to inherit
        # whatever the template's own theme defines, not our brand blue.
        color_type = title_placeholder.text_frame.paragraphs[0].font.color.type
        assert color_type is None

    def test_stat_cards_stay_transparent_so_template_background_shows_through(
        self, template_path
    ):
        pptx_bytes = build_pptx_report(_empty_stats(), template_path=template_path)
        prs = Presentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]

        card_boxes = [
            shape
            for shape in slide.shapes
            if not shape.is_placeholder and shape.has_text_frame and "Total Members" in shape.text_frame.text
        ]
        assert card_boxes
        for box in card_boxes:
            # fill.background() ("no fill", MSO_FILL_TYPE.BACKGROUND) rather
            # than a solid tone colour — lets the template's own background
            # show through instead of painting over it.
            assert box.fill.type.name == "BACKGROUND"

    def test_default_deck_without_a_template_is_unchanged(self):
        # Regression guard: no template at all must still hit the
        # pre-8.33 code path exactly (blank layout, solid tone fills, brand
        # blue text) — Story 8.33 must not touch this path.
        pptx_bytes = build_pptx_report(_empty_stats())
        prs = Presentation(BytesIO(pptx_bytes))
        slide = prs.slides[0]

        assert slide.slide_layout.name == "Blank"
        heading_box = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Members Statistics Report" in shape.text_frame.text
        )
        assert not heading_box.is_placeholder
        assert str(heading_box.text_frame.paragraphs[0].font.color.rgb) == "17458F"

        card_box = next(
            shape
            for shape in slide.shapes
            if shape.has_text_frame and "Total Members" in shape.text_frame.text
        )
        assert card_box.fill.type is not None
