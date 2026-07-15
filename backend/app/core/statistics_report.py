"""Server-side PDF/PPTX export of the members statistics page (Story 2b.14).

Charts are re-rendered from the same MembersStatistics data backing the
live page (via matplotlib) rather than screenshotting the frontend, so the
export stays accurate regardless of the caller's browser/device.
"""

from datetime import date
from io import BytesIO
from pathlib import Path

import matplotlib

matplotlib.use("Agg")  # headless — no display server available on a backend host
import matplotlib.pyplot as plt  # noqa: E402 (must follow matplotlib.use)
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from reportlab.lib import colors  # noqa: E402
from reportlab.lib.pagesizes import letter  # noqa: E402
from reportlab.lib.styles import getSampleStyleSheet  # noqa: E402
from reportlab.lib.units import inch  # noqa: E402
from reportlab.platypus import (  # noqa: E402
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.core.pdf_style import PDF_BODY_FONT_SIZE, PDF_TABLE_HEADER_FONT_SIZE
from app.schemas.member_statistics import MembersStatistics

ROTARY_BLUE = "#17458f"
ROTARY_GOLD = "#f7a81b"
PIE_COLORS = ["#17458f", "#f7a81b", "#5f55ee", "#0f9d9f", "#b3261e", "#9aa4b2"]
CLUB_NAME = "Rotary Club of Discovery Bay"

# Same light/pastel tone pairs as the live page's stat cards
# (frontend/src/index.css --tone-*-bg), so the report visually matches:
# Total/Honorary = blue, New/Countries = lavender, Women/Men = teal,
# Age/Tenure = amber. Two cards share each tone, in display order.
TONE_BLUE_BG = "#e3edfb"
TONE_LAVENDER_BG = "#ece7fb"
TONE_TEAL_BG = "#e0f4f1"
TONE_AMBER_BG = "#fdf0da"
CARD_TONES = [
    TONE_BLUE_BG,
    TONE_BLUE_BG,
    TONE_LAVENDER_BG,
    TONE_LAVENDER_BG,
    TONE_TEAL_BG,
    TONE_TEAL_BG,
    TONE_AMBER_BG,
    TONE_AMBER_BG,
]

# Kept in sync with frontend/src/assets/rotary-logo.png (see that file's own
# comment) — copied here so report generation doesn't depend on the frontend
# source tree being present on whatever host runs the backend.
LOGO_PATH = Path(__file__).resolve().parents[1] / "assets" / "rotary-logo.png"


def _bar_chart_png(labels: list[str], values: list[float], title: str) -> bytes:
    fig, ax = plt.subplots(figsize=(5.0, 2.6))
    ax.bar(labels, values, color=ROTARY_BLUE)
    ax.set_title(title, fontsize=10)
    ax.tick_params(axis="x", rotation=30, labelsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    return _fig_to_png(fig)


def _grouped_bar_chart_png(
    labels: list[str], series: dict[str, list[float]], title: str
) -> bytes:
    fig, ax = plt.subplots(figsize=(5.0, 2.6))
    bar_width = 0.8 / max(len(series), 1)
    x_positions = range(len(labels))
    colors_cycle = [ROTARY_BLUE, "#b3261e", ROTARY_GOLD]
    for index, (name, values) in enumerate(series.items()):
        offsets = [x + index * bar_width for x in x_positions]
        ax.bar(offsets, values, width=bar_width, label=name, color=colors_cycle[index % 3])
    ax.set_xticks([x + bar_width * (len(series) - 1) / 2 for x in x_positions])
    ax.set_xticklabels(labels, rotation=30, fontsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.set_title(title, fontsize=10)
    ax.legend(fontsize=7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    return _fig_to_png(fig)


def _horizontal_bar_chart_png(labels: list[str], values: list[float], title: str) -> bytes:
    """Story 8.32 — mirrors the live pages' `<BarChart layout="vertical">`
    (NGO's Top Organisations / By Classification charts)."""
    fig, ax = plt.subplots(figsize=(5.0, max(2.6, 0.3 * len(labels))))
    ax.barh(labels, values, color=ROTARY_BLUE)
    ax.invert_yaxis()
    ax.set_title(title, fontsize=10)
    ax.tick_params(axis="both", labelsize=7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    return _fig_to_png(fig)


def _line_chart_png(labels: list[str], values: list[float], title: str) -> bytes:
    """Story 8.32 — mirrors the live NGO Statistics page's year-over-year
    `<LineChart>`."""
    fig, ax = plt.subplots(figsize=(5.0, 2.6))
    ax.plot(labels, values, color=ROTARY_GOLD, marker="o", markersize=3)
    ax.set_title(title, fontsize=10)
    ax.tick_params(axis="x", rotation=30, labelsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    return _fig_to_png(fig)


def _pie_chart_png(labels: list[str], values: list[float], title: str) -> bytes:
    fig, ax = plt.subplots(figsize=(5.0, 2.8))
    if sum(values) == 0:
        ax.text(0.5, 0.5, "No data", ha="center", va="center")
        ax.axis("off")
    else:
        ax.pie(
            values,
            labels=labels,
            autopct="%1.0f%%",
            colors=[PIE_COLORS[i % len(PIE_COLORS)] for i in range(len(labels))],
            textprops={"fontsize": 7},
        )
    ax.set_title(title, fontsize=10)
    fig.tight_layout()
    return _fig_to_png(fig)


def _fig_to_png(fig) -> bytes:
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    return buf.getvalue()


def render_charts(stats: MembersStatistics) -> dict[str, bytes]:
    """Returns {chart title: PNG bytes} for every chart shown on the live page."""
    charts: dict[str, bytes] = {}

    charts["Members by join year"] = _bar_chart_png(
        [entry.label for entry in stats.by_join_year],
        [entry.value for entry in stats.by_join_year],
        "Members by join year",
    )

    charts["Growth by Rotary year (joins vs leaves)"] = _grouped_bar_chart_png(
        [entry.label for entry in stats.growth_by_rotary_year],
        {
            "Joins": [entry.joins for entry in stats.growth_by_rotary_year],
            "Leaves": [entry.leaves for entry in stats.growth_by_rotary_year],
        },
        "Growth by Rotary year (joins vs leaves)",
    )

    charts["Nationality distribution"] = _pie_chart_png(
        [entry.label for entry in stats.by_nationality],
        [entry.value for entry in stats.by_nationality],
        "Nationality distribution",
    )

    charts["Tenure distribution (years as Rotarian)"] = _bar_chart_png(
        [entry.label for entry in stats.tenure_distribution],
        [entry.value for entry in stats.tenure_distribution],
        "Tenure distribution (years as Rotarian)",
    )

    charts["Gender distribution"] = _pie_chart_png(
        [entry.label for entry in stats.by_gender],
        [entry.value for entry in stats.by_gender],
        "Gender distribution",
    )

    charts["Age distribution"] = _bar_chart_png(
        [entry.label for entry in stats.age_distribution],
        [entry.value for entry in stats.age_distribution],
        "Age distribution",
    )

    return charts


def _stat_cards() -> list[tuple[str, str]]:
    """(label, attribute) pairs, in the same order/grouping as the live page."""
    return [
        ("Total Members", "total_members"),
        ("Honorary Members", "honorary_members"),
        ("New Members (this Rotary year)", "new_members_this_rotary_year"),
        ("Countries Represented", "countries_represented"),
        ("Number of Women", "women_count"),
        ("Number of Men", "men_count"),
        ("Average Age", "average_age"),
        ("Average Tenure (as Rotarian)", "average_tenure_as_rotarian"),
    ]


def _detail_tables(stats: MembersStatistics) -> list[tuple[str, list[str], list[list[str]]]]:
    """Story 8.13: an "Integral" report adds a detail section restating every
    chart's underlying numbers as a table, on top of everything Simplified
    already shows — using data already on the MembersStatistics response, no
    new backend queries. (8.13 wasn't originally scoped for the Members page —
    this is this story's own reasonable interpretation of "stats + graph +
    detail section" for it, reusing the existing 6 chart datasets.)"""
    return [
        (
            "Members by join year",
            ["Year", "Members"],
            [[entry.label, str(entry.value)] for entry in stats.by_join_year],
        ),
        (
            "Growth by Rotary year",
            ["Rotary year", "Joins", "Leaves"],
            [
                [entry.label, str(entry.joins), str(entry.leaves)]
                for entry in stats.growth_by_rotary_year
            ],
        ),
        (
            "Nationality distribution",
            ["Nationality", "Members"],
            [[entry.label, str(entry.value)] for entry in stats.by_nationality],
        ),
        (
            "Tenure distribution (years as Rotarian)",
            ["Years as Rotarian", "Members"],
            [[entry.label, str(entry.value)] for entry in stats.tenure_distribution],
        ),
        (
            "Gender distribution",
            ["Gender", "Members"],
            [[entry.label, str(entry.value)] for entry in stats.by_gender],
        ),
        (
            "Age distribution",
            ["Age bracket", "Members"],
            [[entry.label, str(entry.value)] for entry in stats.age_distribution],
        ),
    ]


def build_pdf_report(stats: MembersStatistics, report_type: str = "simplified") -> bytes:
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=letter, topMargin=0.5 * inch, bottomMargin=0.5 * inch
    )
    styles = getSampleStyleSheet()
    chart_title_style = styles["Heading4"]
    chart_title_style.fontSize = 9
    chart_title_style.spaceAfter = 2
    story = []

    if LOGO_PATH.exists():
        story.append(Image(str(LOGO_PATH), width=0.6 * inch, height=0.6 * inch))
    story.append(Paragraph(CLUB_NAME, styles["Title"]))
    story.append(Paragraph("Members Statistics Report", styles["Heading2"]))
    story.append(Paragraph(f"Generated {date.today().isoformat()}", styles["Normal"]))
    story.append(Spacer(1, 0.15 * inch))

    card_rows = _stat_cards()
    table_data = []
    for i in range(0, len(card_rows), 4):
        row = card_rows[i : i + 4]
        table_data.append([label for label, _ in row])
        table_data.append([str(getattr(stats, attr) if getattr(stats, attr) is not None else "–") for _, attr in row])
    cards_table = Table(table_data, colWidths=[1.6 * inch] * 4)
    card_style_commands = [
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("FONTSIZE", (0, 1), (-1, 1), 14),
        ("FONTSIZE", (0, 3), (-1, 3), 14),
        ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
        ("FONTNAME", (0, 3), (-1, 3), "Helvetica-Bold"),
        ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor(ROTARY_BLUE)),
        ("TEXTCOLOR", (0, 3), (-1, 3), colors.HexColor(ROTARY_BLUE)),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
    ]
    # Color each card's label+value cell pair with the same tone used on the
    # live page, so the report visually matches (Total/Honorary = blue,
    # New/Countries = lavender, Women/Men = teal, Age/Tenure = amber).
    for card_index in range(len(card_rows)):
        col = card_index % 4
        label_row = 0 if card_index < 4 else 2
        value_row = label_row + 1
        tone = colors.HexColor(CARD_TONES[card_index])
        card_style_commands.append(("BACKGROUND", (col, label_row), (col, value_row), tone))
    cards_table.setStyle(TableStyle(card_style_commands))
    story.append(cards_table)
    story.append(Spacer(1, 0.2 * inch))

    # 2-column chart grid (rather than one full-width chart per section) so
    # all 6 charts plus the header/cards above fit in a 2-page PDF.
    chart_items = list(render_charts(stats).items())
    chart_cell_width, chart_cell_height = 3.3 * inch, 1.9 * inch
    grid_rows = []
    for i in range(0, len(chart_items), 2):
        pair = chart_items[i : i + 2]
        cells = []
        for title, png_bytes in pair:
            cells.append(
                [
                    Paragraph(title, chart_title_style),
                    Image(BytesIO(png_bytes), width=chart_cell_width, height=chart_cell_height),
                ]
            )
        if len(cells) == 1:
            cells.append("")
        grid_rows.append(cells)

    charts_table = Table(grid_rows, colWidths=[3.6 * inch, 3.6 * inch])
    charts_table.setStyle(
        TableStyle(
            [
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    story.append(charts_table)

    if report_type == "integral":
        story.append(PageBreak())
        story.append(Paragraph("Detail — underlying figures", styles["Heading2"]))
        story.append(Spacer(1, 0.1 * inch))
        for title, headers, rows in _detail_tables(stats):
            story.append(Paragraph(title, styles["Heading4"]))
            detail_table = Table([headers] + rows, hAlign="LEFT")
            detail_table.setStyle(
                TableStyle(
                    [
                        ("FONTSIZE", (0, 0), (-1, -1), PDF_BODY_FONT_SIZE),
                        ("FONTSIZE", (0, 0), (-1, 0), PDF_TABLE_HEADER_FONT_SIZE),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(TONE_BLUE_BG)),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                    ]
                )
            )
            story.append(detail_table)
            story.append(Spacer(1, 0.15 * inch))

    doc.build(story)
    return buf.getvalue()


def _pick_blank_layout(prs: Presentation):
    """Story 8.23: an uploaded annual template won't necessarily have a blank
    layout at index 6 like the app's own from-scratch deck does — look for
    one named "blank" first, else fall back to the last available layout
    (emptier than the first, which is usually a title slide)."""
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def build_pptx_report(
    stats: MembersStatistics,
    report_type: str = "simplified",
    template_path: Path | None = None,
) -> bytes:
    # Everything (heading, stat cards, all 6 charts) condensed onto a single
    # widescreen slide rather than one-section-per-slide, per request.
    if template_path is not None:
        prs = Presentation(str(template_path))
        blank_layout = _pick_blank_layout(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

    # The from-scratch layout below is designed against a 13.333in-wide
    # widescreen deck. An uploaded template may use a different slide size
    # (e.g. 10in widescreen, or 4:3) — scale every position/size by the ratio
    # so content roughly fits instead of overflowing or bunching in a corner.
    # `scale` is exactly 1 for the app's own default deck (no-op).
    scale = prs.slide_width / Inches(13.333)

    def sc(emu):
        return int(emu * scale)

    slide = prs.slides.add_slide(blank_layout)

    # The heading box starts right of the logo's actual rendered width (not a
    # hardcoded guess) — the logo's aspect ratio makes its width at a fixed
    # 0.6in height wider than the 0.8in gap a hardcoded left offset assumed,
    # which is what caused the heading text to overlap the logo.
    logo_right_edge = sc(Inches(0.3))
    if LOGO_PATH.exists():
        logo_pic = slide.shapes.add_picture(
            str(LOGO_PATH), sc(Inches(0.3)), sc(Inches(0.2)), height=sc(Inches(0.6))
        )
        logo_right_edge = logo_pic.left + logo_pic.width

    heading_left = logo_right_edge + sc(Inches(0.2))
    heading_box = slide.shapes.add_textbox(
        heading_left,
        sc(Inches(0.18)),
        prs.slide_width - heading_left - sc(Inches(0.3)),
        sc(Inches(0.65)),
    )
    heading_tf = heading_box.text_frame
    heading_tf.text = f"{CLUB_NAME} — Members Statistics Report"
    heading_tf.paragraphs[0].font.size = Pt(22)
    heading_tf.paragraphs[0].font.bold = True
    heading_tf.paragraphs[0].font.color.rgb = RGBColor.from_string("17458F")
    date_p = heading_tf.add_paragraph()
    date_p.text = f"Generated {date.today().isoformat()}"
    date_p.font.size = Pt(11)

    # Stat cards: 4 columns x 2 rows, compact
    card_rows = _stat_cards()
    columns = 4
    card_width, card_height = sc(Inches(2.95)), sc(Inches(0.75))
    gap_x, gap_y = sc(Inches(0.12)), sc(Inches(0.08))
    start_x, start_y = sc(Inches(0.3)), sc(Inches(0.95))
    for index, (label, attr) in enumerate(card_rows):
        col = index % columns
        row = index // columns
        left = start_x + col * (card_width + gap_x)
        top = start_y + row * (card_height + gap_y)
        box = slide.shapes.add_textbox(left, top, card_width, card_height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string(CARD_TONES[index].lstrip("#").upper())
        box.line.fill.background()
        tf = box.text_frame
        tf.margin_left = Pt(6)
        tf.margin_top = Pt(4)
        value = getattr(stats, attr)
        tf.text = str(value if value is not None else "–")
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor.from_string("17458F")
        label_p = tf.add_paragraph()
        label_p.text = label
        label_p.font.size = Pt(9)

    # Charts: 3 columns x 2 rows
    chart_columns = 3
    chart_width = sc(Inches(4.15))
    chart_gap_x, chart_gap_y = sc(Inches(0.08)), sc(Inches(0.15))
    chart_start_x, chart_start_y = sc(Inches(0.25)), sc(Inches(2.75))
    chart_row_height = sc(Inches(2.2))
    for index, (title, png_bytes) in enumerate(render_charts(stats).items()):
        col = index % chart_columns
        row = index // chart_columns
        left = chart_start_x + col * (chart_width + chart_gap_x)
        top = chart_start_y + row * (chart_row_height + chart_gap_y)
        slide.shapes.add_picture(BytesIO(png_bytes), left, top, width=chart_width)

    if report_type == "integral":
        _add_pptx_detail_slides(prs, blank_layout, stats, sc)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_pptx_detail_slides(prs: Presentation, blank_layout, stats: MembersStatistics, sc) -> None:
    """Story 8.13: one detail slide per chart dataset, appended after the
    main Simplified slide, for Integral reports. Tables aren't paginated
    across slides if a dataset has many rows (e.g. many distinct join
    years) — a known simplification given this app's small member base."""
    for title, headers, rows in _detail_tables(stats):
        slide = prs.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(
            sc(Inches(0.4)), sc(Inches(0.3)), sc(Inches(9)), sc(Inches(0.6))
        )
        title_tf = title_box.text_frame
        title_tf.text = title
        title_tf.paragraphs[0].font.size = Pt(20)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = RGBColor.from_string("17458F")

        table_rows = len(rows) + 1
        table_shape = slide.shapes.add_table(
            table_rows,
            len(headers),
            sc(Inches(0.4)),
            sc(Inches(1.1)),
            sc(Inches(9)),
            sc(Inches(min(0.4 * table_rows, 5.5))),
        )
        table = table_shape.table
        for col_index, header in enumerate(headers):
            cell = table.cell(0, col_index)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        for row_index, row_values in enumerate(rows, start=1):
            for col_index, value in enumerate(row_values):
                cell = table.cell(row_index, col_index)
                cell.text = value
                cell.text_frame.paragraphs[0].font.size = Pt(10)
