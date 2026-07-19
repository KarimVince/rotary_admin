"""PDF/PPT export of the Summary page (Story 14.10) — cards + the 3
breakdown charts, reusing the same matplotlib chart-drawing conventions as
statistics_report.py (own copies here since the chart shapes differ:
this page only ever needs a pie and a simple bar, not the full set)."""
from datetime import date as date_type
from io import BytesIO

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from reportlab.lib import colors  # noqa: E402
from reportlab.lib.pagesizes import letter  # noqa: E402
from reportlab.lib.styles import getSampleStyleSheet  # noqa: E402
from reportlab.lib.units import inch  # noqa: E402
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle  # noqa: E402

from app.core.dinner_forecast_report import INTL_LOGO_PATH
from app.core.statistics_report import CLUB_NAME, LOGO_PATH, PIE_COLORS, ROTARY_BLUE
from app.schemas.event_summary import EventSummary

PAGE_MARGIN = 0.5 * inch


def _fig_to_png(fig) -> bytes:
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    return buf.getvalue()


def _pie_chart_png(labels: list[str], values: list[float], title: str) -> bytes:
    fig, ax = plt.subplots(figsize=(4.5, 2.8))
    if sum(values) == 0:
        ax.text(0.5, 0.5, "No data", ha="center", va="center")
        ax.axis("off")
    else:
        ax.pie(
            values,
            labels=labels,
            autopct="%1.0f%%",
            colors=[PIE_COLORS[i % len(PIE_COLORS)] for i in range(len(labels))],
            textprops={"fontsize": 7},
        )
    ax.set_title(title, fontsize=10)
    fig.tight_layout()
    return _fig_to_png(fig)


def _bar_chart_png(labels: list[str], values: list[float], title: str) -> bytes:
    fig, ax = plt.subplots(figsize=(4.5, 2.8))
    ax.bar(labels, values, color=ROTARY_BLUE)
    ax.set_title(title, fontsize=10)
    ax.tick_params(axis="x", labelsize=7)
    ax.tick_params(axis="y", labelsize=7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    return _fig_to_png(fig)


def _render_charts(summary: EventSummary) -> dict[str, bytes]:
    return {
        "Revenue breakdown": _pie_chart_png(
            [e.label for e in summary.revenue_breakdown],
            [e.value for e in summary.revenue_breakdown],
            "Revenue breakdown",
        ),
        "Cost breakdown": (
            _pie_chart_png(
                [e.label for e in summary.cost_breakdown],
                [e.value for e in summary.cost_breakdown],
                "Cost breakdown",
            )
            if summary.cost_breakdown
            else _pie_chart_png([], [], "Cost breakdown")
        ),
        "Result overview": _bar_chart_png(
            [e.label for e in summary.result_overview],
            [e.value for e in summary.result_overview],
            "Result overview",
        ),
    }


def _card_rows(summary: EventSummary) -> list[tuple[str, str]]:
    rows = [
        ("Total Raised", f"{summary.total_raised:.2f}"),
        ("Auction Total", f"{summary.auction_total:.2f}"),
        ("Lucky Draw Total", f"{summary.lucky_draw_total:.2f}"),
        ("Other Donation", f"{summary.other_donation:.2f}"),
        ("Total Revenue", f"{summary.total_revenue:.2f}"),
        ("Ticket Revenue", f"{summary.ticket_revenue:.2f}"),
        ("Sponsor Revenue", f"{summary.sponsor_revenue:.2f}"),
        ("Total Cost", f"{summary.total_cost:.2f}"),
    ]
    for entry in summary.cost_per_category:
        rows.append((f"Cost — {entry.label}", f"{entry.value:.2f}"))
    rows.append(("Net Operational Result", f"{summary.net_operational_result:.2f}"))
    return rows


def build_pdf_report(event_name: str, event_date: date_type, summary: EventSummary) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        topMargin=PAGE_MARGIN,
        bottomMargin=PAGE_MARGIN,
        leftMargin=PAGE_MARGIN,
        rightMargin=PAGE_MARGIN,
    )
    styles = getSampleStyleSheet()
    usable_width = letter[0] - 2 * PAGE_MARGIN
    story = []

    left_logo = Image(str(LOGO_PATH), width=0.6 * inch, height=0.6 * inch) if LOGO_PATH.exists() else ""
    right_logo = (
        Image(str(INTL_LOGO_PATH), width=0.6 * inch, height=0.6 * inch)
        if INTL_LOGO_PATH.exists()
        else ""
    )
    if left_logo or right_logo:
        logo_row = Table([[left_logo, right_logo]], colWidths=[usable_width / 2, usable_width / 2])
        logo_row.setStyle(
            TableStyle(
                [
                    ("ALIGN", (0, 0), (0, 0), "LEFT"),
                    ("ALIGN", (1, 0), (1, 0), "RIGHT"),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]
            )
        )
        story.append(logo_row)

    story.append(Paragraph(CLUB_NAME, styles["Title"]))
    story.append(
        Paragraph(f"Event Summary — {event_name} ({event_date.strftime('%d %b %Y')})", styles["Heading2"])
    )
    story.append(Spacer(1, 0.2 * inch))

    header_style = styles["BodyText"].clone("cell")
    header_style.fontName = "Helvetica-Bold"
    header_style.textColor = colors.white
    body_style = styles["BodyText"].clone("cell2")

    rows = [[Paragraph("Metric", header_style), Paragraph("Value (HKD)", header_style)]]
    for label, value in _card_rows(summary):
        rows.append([Paragraph(label, body_style), Paragraph(value, body_style)])

    table = Table(rows, colWidths=[3.5 * inch, 2.0 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(ROTARY_BLUE)),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#d0d5dd")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f7fa")]),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 0.3 * inch))

    for title, png_bytes in _render_charts(summary).items():
        story.append(Image(BytesIO(png_bytes), width=4.5 * inch, height=2.8 * inch))
        story.append(Spacer(1, 0.15 * inch))

    doc.build(story)
    return buffer.getvalue()


def build_pptx_report(event_name: str, event_date: date_type, summary: EventSummary) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank_layout)
    title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = f"{CLUB_NAME} — Event Summary"
    title_frame.paragraphs[0].font.size = Pt(32)
    subtitle_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(1))
    subtitle_box.text_frame.text = f"{event_name} ({event_date.strftime('%d %b %Y')})"

    cards_slide = prs.slides.add_slide(blank_layout)
    rows = _card_rows(summary)
    table_shape = cards_slide.shapes.add_table(
        len(rows) + 1, 2, Inches(0.5), Inches(0.5), Inches(8), Inches(0.4 * (len(rows) + 1))
    )
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value (HKD)"
    for index, (label, value) in enumerate(rows, start=1):
        table.cell(index, 0).text = label
        table.cell(index, 1).text = value

    for title, png_bytes in _render_charts(summary).items():
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6)).text_frame.text = title
        slide.shapes.add_picture(BytesIO(png_bytes), Inches(1.5), Inches(1.0), width=Inches(6))

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
