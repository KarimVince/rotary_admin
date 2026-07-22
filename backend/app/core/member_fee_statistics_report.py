"""Story 16.16 — PDF/PPTX export of the Member Fees Statistics tab.

Mirrors the pattern established for Members (Story 2b.14) and NGO Donations
(Story 8.32): charts re-rendered server-side via matplotlib from the same
data backing the live page, reusing shared constants/chart helpers from
`statistics_report.py`. No report-type variants or PPT template support —
the live tab has neither, so the export is a single simplified view: the
3 stat cards plus the 2 history charts (Amount collected, Paying members).
"""

from datetime import date
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from app.core.statistics_report import (
    CLUB_NAME,
    LOGO_PATH,
    ROTARY_BLUE,
    TONE_BLUE_BG,
    TONE_TEAL_BG,
    _bar_chart_png,
    _grouped_bar_chart_png,
    _pick_blank_layout,
    add_heading,
    style_card_fill,
    style_card_text_color,
)
from app.schemas.member_fee import FeeYearHistory, MemberFeeStatistics

REPORT_TITLE = "Fee Statistics Report"
# Same tone-per-card as the live page's StatCard bg props (MemberFees.jsx):
# average fee = blue, collected = teal, outstanding = rose.
TONE_ROSE_BG = "#fbe3e6"
CARD_TONES = [TONE_BLUE_BG, TONE_TEAL_BG, TONE_ROSE_BG]


def _rotary_year_label(year: int) -> str:
    return f"{year}–{year + 1}"


def _format_currency(value: float, currency: str | None) -> str:
    formatted = f"{value:,.0f}"
    return f"{formatted} {currency}" if currency else formatted


def stat_cards(stats: MemberFeeStatistics) -> list[tuple[str, str]]:
    year_label = _rotary_year_label(stats.rotary_year)
    return [
        (
            f"Average fee per active member — {year_label}",
            _format_currency(stats.average_fee_per_active_member, stats.currency),
        ),
        (
            f"Total collected — {year_label}",
            _format_currency(stats.total_collected, stats.currency),
        ),
        (
            f"Total outstanding — {year_label}",
            _format_currency(stats.total_outstanding, stats.currency),
        ),
    ]


def render_charts(history: list[FeeYearHistory]) -> dict[str, bytes]:
    year_labels = [_rotary_year_label(row.rotary_year) for row in history]
    return {
        "Amount collected over years": _bar_chart_png(
            year_labels,
            [row.total_collected for row in history],
            "Amount collected over years",
        ),
        "Paying members over years": _grouped_bar_chart_png(
            year_labels,
            {
                "Paid": [row.paid_count for row in history],
                "Zero": [row.zero_count for row in history],
            },
            "Paying members over years",
        ),
    }


def build_pdf_report(stats: MemberFeeStatistics, history: list[FeeYearHistory]) -> bytes:
    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.5 * inch, bottomMargin=0.5 * inch)
    styles = getSampleStyleSheet()
    story = []

    if LOGO_PATH.exists():
        story.append(Image(str(LOGO_PATH), width=0.6 * inch, height=0.6 * inch))
    story.append(Paragraph(CLUB_NAME, styles["Title"]))
    story.append(Paragraph(REPORT_TITLE, styles["Heading2"]))
    story.append(
        Paragraph(f"Rotary Year {_rotary_year_label(stats.rotary_year)}", styles["Normal"])
    )
    story.append(Paragraph(f"Generated {date.today().isoformat()}", styles["Normal"]))
    story.append(Spacer(1, 0.15 * inch))

    cards = stat_cards(stats)
    table_data = [[label for label, _ in cards], [value for _, value in cards]]
    cards_table = Table(table_data, colWidths=[2.3 * inch] * 3)
    card_style_commands = [
        ("FONTSIZE", (0, 0), (-1, 0), 8),
        ("FONTSIZE", (0, 1), (-1, 1), 14),
        ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
        ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor(ROTARY_BLUE)),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
    ]
    for col in range(len(cards)):
        tone = colors.HexColor(CARD_TONES[col])
        card_style_commands.append(("BACKGROUND", (col, 0), (col, 1), tone))
    cards_table.setStyle(TableStyle(card_style_commands))
    story.append(cards_table)
    story.append(Spacer(1, 0.2 * inch))

    story.append(
        Paragraph(
            f"{stats.paid_count} paid, {stats.unpaid_count} unpaid, {stats.total_members} total "
            f"member{'s' if stats.total_members != 1 else ''} billed. "
            f"Collection rate {stats.collection_rate:.1f}%.",
            styles["Normal"],
        )
    )
    story.append(Spacer(1, 0.2 * inch))

    if not history:
        story.append(Paragraph("No fee history recorded yet.", styles["Normal"]))
    else:
        chart_title_style = styles["Heading4"]
        chart_title_style.fontSize = 9
        chart_title_style.spaceAfter = 2
        chart_items = list(render_charts(history).items())
        chart_cell_width, chart_cell_height = 3.3 * inch, 2.4 * inch
        row = []
        for title, png_bytes in chart_items:
            row.append(
                [
                    Paragraph(title, chart_title_style),
                    Image(BytesIO(png_bytes), width=chart_cell_width, height=chart_cell_height),
                ]
            )
        charts_table = Table([row], colWidths=[3.6 * inch, 3.6 * inch])
        charts_table.setStyle(
            TableStyle(
                [
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]
            )
        )
        story.append(charts_table)

    doc.build(story)
    return buf.getvalue()


def build_pptx_report(stats: MemberFeeStatistics, history: list[FeeYearHistory]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = _pick_blank_layout(prs)

    slide = prs.slides.add_slide(blank_layout)

    logo_right_edge = Inches(0.3)
    if LOGO_PATH.exists():
        logo_pic = slide.shapes.add_picture(
            str(LOGO_PATH), Inches(0.3), Inches(0.2), height=Inches(0.6)
        )
        logo_right_edge = logo_pic.left + logo_pic.width

    heading_left = logo_right_edge + Inches(0.2)
    add_heading(
        slide,
        False,
        f"{CLUB_NAME} — {REPORT_TITLE}",
        f"Rotary Year {_rotary_year_label(stats.rotary_year)} — Generated {date.today().isoformat()}",
        (heading_left, Inches(0.18), prs.slide_width - heading_left - Inches(0.3), Inches(0.65)),
    )

    cards = stat_cards(stats)
    columns = 3
    card_width, card_height = Inches(4.0), Inches(0.85)
    gap_x = Inches(0.15)
    start_x, start_y = Inches(0.3), Inches(1.0)
    for index, (label, value) in enumerate(cards):
        left = start_x + index * (card_width + gap_x)
        box = slide.shapes.add_textbox(left, start_y, card_width, card_height)
        style_card_fill(box, CARD_TONES[index % columns], False)
        tf = box.text_frame
        tf.margin_left = Pt(6)
        tf.margin_top = Pt(4)
        tf.text = value
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        style_card_text_color(tf.paragraphs[0], False)
        label_p = tf.add_paragraph()
        label_p.text = label
        label_p.font.size = Pt(9)

    summary_box = slide.shapes.add_textbox(start_x, Inches(2.05), Inches(12.5), Inches(0.4))
    summary_box.text_frame.text = (
        f"{stats.paid_count} paid, {stats.unpaid_count} unpaid, {stats.total_members} total "
        f"member{'s' if stats.total_members != 1 else ''} billed. "
        f"Collection rate {stats.collection_rate:.1f}%."
    )
    summary_box.text_frame.paragraphs[0].font.size = Pt(12)

    if history:
        chart_width = Inches(6.0)
        chart_start_x, chart_start_y = Inches(0.4), Inches(2.6)
        gap_x = Inches(0.3)
        for index, (title, png_bytes) in enumerate(render_charts(history).items()):
            left = chart_start_x + index * (chart_width + gap_x)
            slide.shapes.add_picture(BytesIO(png_bytes), left, chart_start_y, width=chart_width)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
