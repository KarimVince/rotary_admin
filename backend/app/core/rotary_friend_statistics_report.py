"""Story 8.32 — PDF/PPTX export of the Rotary Friends statistics page.

Mirrors the pattern established for Members (Story 2b.14) and NGO/Donations
(Story 8.32) reports: charts re-rendered server-side via matplotlib from the
same data backing the live page, reusing the shared constants/chart helpers
from `statistics_report.py` rather than a new generator.

Scope note (flagged in ClickUp on this story): the "Integral" detail section
uses each Friend's name, contact channel, source, and tags — the only
per-friend attributes that actually exist on the `RotaryFriend` model. There
is no per-year "engagement metrics" data anywhere in this app (no event
link, no activity log, no date-bearing field besides created_at/updated_at),
so the story's literal "relevant engagement metrics for the selected year"
isn't buildable as specified — this report has no year concept at all,
matching the live page (which has no year filter either).
"""

from datetime import date
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.core.pdf_style import PDF_BODY_FONT_SIZE, PDF_TABLE_HEADER_FONT_SIZE
from app.core.statistics_report import (
    CLUB_NAME,
    LOGO_PATH,
    ROTARY_BLUE,
    TONE_BLUE_BG,
    _bar_chart_png,
    _pick_blank_layout,
    _pie_chart_png,
    _title_placeholder,
    add_heading,
    style_card_fill,
    style_card_text_color,
)
from app.models import RotaryFriend
from app.schemas.rotary_friend_statistics import RotaryFriendStatistics

REPORT_TITLE = "Friends of Rotary Statistics"


def render_charts(stats: RotaryFriendStatistics) -> dict[str, bytes]:
    charts: dict[str, bytes] = {}
    charts["By source"] = _bar_chart_png(
        [entry.label for entry in stats.by_source],
        [entry.value for entry in stats.by_source],
        "By source",
    )
    charts["By tag"] = _bar_chart_png(
        [entry.label for entry in stats.by_tag],
        [entry.value for entry in stats.by_tag],
        "By tag",
    )
    charts["Contactability"] = _pie_chart_png(
        [entry.label for entry in stats.contactability],
        [entry.value for entry in stats.contactability],
        "Contactability",
    )
    return charts


def _friend_detail_rows(friends: list[RotaryFriend]) -> list[list[str]]:
    rows = []
    for friend in friends:
        name = f"{friend.first_name} {friend.last_name}".strip()
        contact = ", ".join(
            filter(None, [friend.email and "Email", friend.whatsapp and "WhatsApp"])
        ) or "—"
        rows.append([name, contact, friend.source or "—", friend.tags or "—"])
    return rows


def build_pdf_report(
    stats: RotaryFriendStatistics, friends: list[RotaryFriend], report_type: str = "simplified"
) -> bytes:
    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.5 * inch, bottomMargin=0.5 * inch)
    styles = getSampleStyleSheet()
    story = []

    report_type_label = "Integral" if report_type == "integral" else "Simplified"

    if LOGO_PATH.exists():
        story.append(Image(str(LOGO_PATH), width=0.6 * inch, height=0.6 * inch))
    story.append(Paragraph(CLUB_NAME, styles["Title"]))
    story.append(Paragraph(f"{REPORT_TITLE} — {report_type_label} Report", styles["Heading2"]))
    story.append(Paragraph(f"Generated {date.today().isoformat()}", styles["Normal"]))
    story.append(Spacer(1, 0.15 * inch))

    card_table = Table(
        [["Total Friends"], [str(stats.total_friends)]], colWidths=[2.0 * inch]
    )
    card_table.setStyle(
        TableStyle(
            [
                ("FONTSIZE", (0, 0), (-1, 0), 8),
                ("FONTSIZE", (0, 1), (-1, 1), 18),
                ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
                ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor(ROTARY_BLUE)),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(TONE_BLUE_BG)),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
            ]
        )
    )
    story.append(card_table)
    story.append(Spacer(1, 0.2 * inch))

    chart_title_style = styles["Heading4"]
    chart_title_style.fontSize = 9
    chart_title_style.spaceAfter = 2
    for title, png_bytes in render_charts(stats).items():
        story.append(Paragraph(title, chart_title_style))
        story.append(Image(BytesIO(png_bytes), width=6.5 * inch, height=3.0 * inch))
        story.append(Spacer(1, 0.15 * inch))

    if report_type == "integral":
        story.append(PageBreak())
        story.append(Paragraph("Detail — Friends", styles["Heading2"]))
        story.append(Spacer(1, 0.1 * inch))
        headers = ["Name", "Contact", "Source", "Tags"]
        rows = _friend_detail_rows(friends)
        if rows:
            detail_table = Table([headers] + rows, hAlign="LEFT")
            detail_table.setStyle(
                TableStyle(
                    [
                        ("FONTSIZE", (0, 0), (-1, -1), PDF_BODY_FONT_SIZE),
                        ("FONTSIZE", (0, 0), (-1, 0), PDF_TABLE_HEADER_FONT_SIZE),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(TONE_BLUE_BG)),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                    ]
                )
            )
            story.append(detail_table)
        else:
            story.append(Paragraph("No Rotary Friends recorded yet.", styles["Normal"]))

    doc.build(story)
    return buf.getvalue()


def build_pptx_report(
    stats: RotaryFriendStatistics,
    friends: list[RotaryFriend],
    report_type: str = "simplified",
    template_path: BytesIO | None = None,
) -> bytes:
    report_type_label = "Integral" if report_type == "integral" else "Simplified"

    using_template = template_path is not None
    if using_template:
        prs = Presentation(template_path)
        blank_layout = _pick_blank_layout(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

    scale = prs.slide_width / Inches(13.333)

    def sc(emu):
        return int(emu * scale)

    slide = prs.slides.add_slide(blank_layout)

    logo_right_edge = sc(Inches(0.3))
    if LOGO_PATH.exists():
        logo_pic = slide.shapes.add_picture(
            str(LOGO_PATH), sc(Inches(0.3)), sc(Inches(0.2)), height=sc(Inches(0.6))
        )
        logo_right_edge = logo_pic.left + logo_pic.width

    heading_left = logo_right_edge + sc(Inches(0.2))
    add_heading(
        slide,
        using_template,
        f"{CLUB_NAME} — {REPORT_TITLE} — {report_type_label} Report",
        f"Generated {date.today().isoformat()}",
        (heading_left, sc(Inches(0.18)), prs.slide_width - heading_left - sc(Inches(0.3)), sc(Inches(0.65))),
    )

    card_box = slide.shapes.add_textbox(sc(Inches(0.3)), sc(Inches(0.95)), sc(Inches(2.5)), sc(Inches(0.9)))
    style_card_fill(card_box, TONE_BLUE_BG, using_template)
    card_tf = card_box.text_frame
    card_tf.margin_left = Pt(8)
    card_tf.margin_top = Pt(6)
    card_tf.text = str(stats.total_friends)
    card_tf.paragraphs[0].font.size = Pt(24)
    card_tf.paragraphs[0].font.bold = True
    style_card_text_color(card_tf.paragraphs[0], using_template)
    label_p = card_tf.add_paragraph()
    label_p.text = "Total Friends"
    label_p.font.size = Pt(10)

    chart_columns = 3
    chart_width = sc(Inches(4.15))
    chart_gap_x = sc(Inches(0.1))
    chart_start_x, chart_start_y = sc(Inches(0.3)), sc(Inches(2.15))
    for index, (title, png_bytes) in enumerate(render_charts(stats).items()):
        left = chart_start_x + index % chart_columns * (chart_width + chart_gap_x)
        slide.shapes.add_picture(BytesIO(png_bytes), left, chart_start_y, width=chart_width)

    if report_type == "integral":
        _add_detail_slide(prs, blank_layout, friends, sc, using_template)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_detail_slide(
    prs: Presentation, blank_layout, friends: list[RotaryFriend], sc, using_template: bool = False
) -> None:
    slide = prs.slides.add_slide(blank_layout)

    title_text = "Detail — Friends"
    placeholder = _title_placeholder(slide) if using_template else None
    if placeholder is not None:
        placeholder.text_frame.text = title_text
    else:
        title_box = slide.shapes.add_textbox(
            sc(Inches(0.4)), sc(Inches(0.3)), sc(Inches(9)), sc(Inches(0.6))
        )
        title_tf = title_box.text_frame
        title_tf.text = title_text
        title_tf.paragraphs[0].font.size = Pt(20)
        title_tf.paragraphs[0].font.bold = True
        if not using_template:
            title_tf.paragraphs[0].font.color.rgb = RGBColor.from_string("17458F")

    headers = ["Name", "Contact", "Source", "Tags"]
    rows = _friend_detail_rows(friends)
    if not rows:
        empty_box = slide.shapes.add_textbox(
            sc(Inches(0.4)), sc(Inches(1.1)), sc(Inches(9)), sc(Inches(0.5))
        )
        empty_box.text_frame.text = "No Rotary Friends recorded yet."
        return

    table_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(
        table_rows, len(headers), sc(Inches(0.4)), sc(Inches(1.1)), sc(Inches(9)),
        sc(Inches(min(0.4 * table_rows, 5.5))),
    )
    table = table_shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for row_index, row_values in enumerate(rows, start=1):
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(10)
