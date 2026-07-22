"""Story 8.32 — PDF/PPTX export of the NGO (Donations) statistics page.

Mirrors the pattern established for Members (Story 2b.14): charts
re-rendered server-side via matplotlib from the same data backing the live
page, reusing shared constants/chart helpers from `statistics_report.py`.

Donation totals are per-currency and never summed across currencies (Story
3.7), same as the live page — a report covers exactly one currency at a
time, chosen by the caller (defaults to the first currency block, same as
the frontend's initial selection).
"""

from datetime import date
from io import BytesIO
from pathlib import Path

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.core.config import settings
from app.core.statistics_report import (
    CLUB_NAME,
    LOGO_PATH,
    ROTARY_BLUE,
    TONE_AMBER_BG,
    TONE_BLUE_BG,
    TONE_LAVENDER_BG,
    TONE_TEAL_BG,
    _bar_chart_png,
    _horizontal_bar_chart_png,
    _line_chart_png,
    _pick_blank_layout,
    _title_placeholder,
    add_heading,
    style_card_fill,
    style_card_text_color,
)
from app.schemas.donation_statistics import CurrencyStatistics, DonationStatistics

REPORT_TITLE = "NGO & Services Project — Statistics"
TOP_ORGS_LIMIT = 10
CARD_TONES = [TONE_BLUE_BG, TONE_BLUE_BG, TONE_LAVENDER_BG, TONE_TEAL_BG, TONE_TEAL_BG, TONE_AMBER_BG]


def _format_currency(value: float, currency: str) -> str:
    return f"{value:,.0f} {currency}"


def _rotary_year_label(year: int) -> str:
    return f"{year}–{year + 1}"


def _current_currency_stats(
    stats: DonationStatistics, currency: str | None
) -> CurrencyStatistics | None:
    if not stats.by_currency:
        return None
    if currency is None:
        return stats.by_currency[0]
    return next((block for block in stats.by_currency if block.currency == currency), None)


def stat_cards(stats: DonationStatistics) -> list[tuple[str, str]]:
    year_label = _rotary_year_label(stats.selected_rotary_year)
    return [
        ("Total donated (all-time)", _format_currency(stats.all_time.total_hkd, "HKD")),
        ("Total donated (all-time)", _format_currency(stats.all_time.total_usd, "USD")),
        ("Organisations supported (all-time)", str(stats.all_time_organisations_count)),
        (f"Total donated — {year_label}", _format_currency(stats.selected_year.total_hkd, "HKD")),
        (f"Total donated — {year_label}", _format_currency(stats.selected_year.total_usd, "USD")),
        (f"Organisations supported — {year_label}", str(stats.selected_year_organisations_count)),
    ]


def render_charts(
    stats: DonationStatistics, currency_stats: CurrencyStatistics | None
) -> dict[str, bytes]:
    charts: dict[str, bytes] = {}
    if currency_stats is None:
        return charts

    year_labels = [entry.label for entry in currency_stats.total_by_rotary_year]
    year_values = [entry.value for entry in currency_stats.total_by_rotary_year]
    charts[f"Total donated per rotary year ({currency_stats.currency})"] = _bar_chart_png(
        year_labels, year_values, "Total donated per rotary year"
    )
    charts[f"Year-over-year trend ({currency_stats.currency})"] = _line_chart_png(
        year_labels, year_values, "Year-over-year trend"
    )

    top_orgs_selected = currency_stats.total_by_organisation_selected_year[:TOP_ORGS_LIMIT]
    charts["Top organisations — Selected Year"] = _horizontal_bar_chart_png(
        [entry.label for entry in top_orgs_selected],
        [entry.value for entry in top_orgs_selected],
        "Top organisations — Selected Year",
    )
    charts["By classification — Selected Year"] = _horizontal_bar_chart_png(
        [entry.label for entry in currency_stats.total_by_classification],
        [entry.value for entry in currency_stats.total_by_classification],
        "By classification — Selected Year",
    )

    top_orgs_all_time = currency_stats.total_by_organisation[:TOP_ORGS_LIMIT]
    charts["Top organisations — All Years"] = _horizontal_bar_chart_png(
        [entry.label for entry in top_orgs_all_time],
        [entry.value for entry in top_orgs_all_time],
        "Top organisations — All Years",
    )
    charts["By classification — All Years"] = _horizontal_bar_chart_png(
        [entry.label for entry in currency_stats.total_by_classification_all_time],
        [entry.value for entry in currency_stats.total_by_classification_all_time],
        "By classification — All Years",
    )
    return charts


def resolve_logo_bytes(logo_url: str | None) -> BytesIO | None:
    """Story 16.6 — logo_url is now a full Supabase Storage public URL for any
    logo uploaded after the migration; fetch it directly. A pre-migration
    relative "/static/..." path is tried against local disk as a best-effort
    fallback, in case the ephemeral filesystem still happens to have it since
    the last restart."""
    if not logo_url:
        return None
    if logo_url.startswith("http://") or logo_url.startswith("https://"):
        try:
            response = httpx.get(logo_url, timeout=10.0)
        except httpx.HTTPError:
            return None
        return BytesIO(response.content) if response.status_code == 200 else None
    filename = logo_url.rsplit("/", 1)[-1]
    path = Path(settings.upload_dir) / "organisations" / filename
    return BytesIO(path.read_bytes()) if path.exists() else None


def build_pdf_report(
    stats: DonationStatistics,
    currency: str | None,
    ngo_breakdown: list[dict],
    report_type: str = "simplified",
) -> bytes:
    currency_stats = _current_currency_stats(stats, currency)
    report_type_label = "Integral" if report_type == "integral" else "Simplified"

    buf = BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.5 * inch, bottomMargin=0.5 * inch)
    styles = getSampleStyleSheet()
    story = []

    if LOGO_PATH.exists():
        story.append(Image(str(LOGO_PATH), width=0.6 * inch, height=0.6 * inch))
    story.append(Paragraph(CLUB_NAME, styles["Title"]))
    story.append(Paragraph(f"{REPORT_TITLE} — {report_type_label} Report", styles["Heading2"]))
    story.append(
        Paragraph(f"Rotary Year {_rotary_year_label(stats.selected_rotary_year)}", styles["Normal"])
    )
    story.append(Paragraph(f"Generated {date.today().isoformat()}", styles["Normal"]))
    story.append(Spacer(1, 0.15 * inch))

    cards = stat_cards(stats)
    table_data = [[label for label, _ in cards[:3]], [value for _, value in cards[:3]]]
    table_data += [[label for label, _ in cards[3:]], [value for _, value in cards[3:]]]
    cards_table = Table(table_data, colWidths=[2.1 * inch] * 3)
    card_style_commands = [
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("FONTSIZE", (0, 1), (-1, 1), 13),
        ("FONTSIZE", (0, 3), (-1, 3), 13),
        ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
        ("FONTNAME", (0, 3), (-1, 3), "Helvetica-Bold"),
        ("TEXTCOLOR", (0, 1), (-1, 1), colors.HexColor(ROTARY_BLUE)),
        ("TEXTCOLOR", (0, 3), (-1, 3), colors.HexColor(ROTARY_BLUE)),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
    ]
    for card_index in range(len(cards)):
        col = card_index % 3
        label_row = 0 if card_index < 3 else 2
        value_row = label_row + 1
        tone = colors.HexColor(CARD_TONES[card_index])
        card_style_commands.append(("BACKGROUND", (col, label_row), (col, value_row), tone))
    cards_table.setStyle(TableStyle(card_style_commands))
    story.append(cards_table)
    story.append(Spacer(1, 0.2 * inch))

    if currency_stats is None:
        story.append(Paragraph("No donations recorded yet.", styles["Normal"]))
    else:
        chart_title_style = styles["Heading4"]
        chart_title_style.fontSize = 9
        chart_title_style.spaceAfter = 2
        chart_items = list(render_charts(stats, currency_stats).items())
        chart_cell_width, chart_cell_height = 3.3 * inch, 1.9 * inch
        grid_rows = []
        for i in range(0, len(chart_items), 2):
            pair = chart_items[i : i + 2]
            cells = []
            for title, png_bytes in pair:
                cells.append(
                    [
                        Paragraph(title, chart_title_style),
                        Image(BytesIO(png_bytes), width=chart_cell_width, height=chart_cell_height),
                    ]
                )
            if len(cells) == 1:
                cells.append("")
            grid_rows.append(cells)
        charts_table = Table(grid_rows, colWidths=[3.6 * inch, 3.6 * inch])
        charts_table.setStyle(
            TableStyle(
                [
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]
            )
        )
        story.append(charts_table)

    if report_type == "integral":
        story.append(PageBreak())
        story.append(
            Paragraph(
                f"Detail — Organisations funded ({_rotary_year_label(stats.selected_rotary_year)})",
                styles["Heading2"],
            )
        )
        story.append(Spacer(1, 0.1 * inch))
        if not ngo_breakdown:
            story.append(Paragraph("No organisations funded this year.", styles["Normal"]))
        else:
            rows = []
            for row in ngo_breakdown:
                logo_bytes = row.get("logo_bytes")
                logo_cell = Image(logo_bytes, width=0.3 * inch, height=0.3 * inch) if logo_bytes else "—"
                rows.append(
                    [logo_cell, row["name"], _format_currency(row["total"], currency_stats.currency)]
                )
            detail_table = Table(
                [["", "Organisation", "Total donated"]] + rows,
                colWidths=[0.45 * inch, 3.5 * inch, 2.0 * inch],
                hAlign="LEFT",
            )
            detail_table.setStyle(
                TableStyle(
                    [
                        ("FONTSIZE", (0, 0), (-1, -1), 9),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dde3ec")),
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(TONE_BLUE_BG)),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                    ]
                )
            )
            story.append(detail_table)

    doc.build(story)
    return buf.getvalue()


def build_pptx_report(
    stats: DonationStatistics,
    currency: str | None,
    ngo_breakdown: list[dict],
    report_type: str = "simplified",
    template_path: BytesIO | None = None,
) -> bytes:
    currency_stats = _current_currency_stats(stats, currency)
    report_type_label = "Integral" if report_type == "integral" else "Simplified"

    using_template = template_path is not None
    if using_template:
        prs = Presentation(template_path)
        blank_layout = _pick_blank_layout(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

    scale = prs.slide_width / Inches(13.333)

    def sc(emu):
        return int(emu * scale)

    slide = prs.slides.add_slide(blank_layout)

    logo_right_edge = sc(Inches(0.3))
    if LOGO_PATH.exists():
        logo_pic = slide.shapes.add_picture(
            str(LOGO_PATH), sc(Inches(0.3)), sc(Inches(0.2)), height=sc(Inches(0.6))
        )
        logo_right_edge = logo_pic.left + logo_pic.width

    heading_left = logo_right_edge + sc(Inches(0.2))
    add_heading(
        slide,
        using_template,
        f"{CLUB_NAME} — {REPORT_TITLE} — {report_type_label} Report",
        f"Rotary Year {_rotary_year_label(stats.selected_rotary_year)} — "
        f"Generated {date.today().isoformat()}",
        (heading_left, sc(Inches(0.18)), prs.slide_width - heading_left - sc(Inches(0.3)), sc(Inches(0.65))),
    )

    cards = stat_cards(stats)
    columns = 3
    card_width, card_height = sc(Inches(3.9)), sc(Inches(0.7))
    gap_x, gap_y = sc(Inches(0.12)), sc(Inches(0.08))
    start_x, start_y = sc(Inches(0.3)), sc(Inches(0.95))
    for index, (label, value) in enumerate(cards):
        col = index % columns
        row = index // columns
        left = start_x + col * (card_width + gap_x)
        top = start_y + row * (card_height + gap_y)
        box = slide.shapes.add_textbox(left, top, card_width, card_height)
        style_card_fill(box, CARD_TONES[index], using_template)
        tf = box.text_frame
        tf.margin_left = Pt(6)
        tf.margin_top = Pt(4)
        tf.text = value
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        style_card_text_color(tf.paragraphs[0], using_template)
        label_p = tf.add_paragraph()
        label_p.text = label
        label_p.font.size = Pt(9)

    if currency_stats is not None:
        chart_columns = 3
        chart_width = sc(Inches(4.15))
        chart_gap_x, chart_gap_y = sc(Inches(0.08)), sc(Inches(0.15))
        chart_start_x, chart_start_y = sc(Inches(0.25)), sc(Inches(2.75))
        chart_row_height = sc(Inches(2.2))
        for index, (title, png_bytes) in enumerate(render_charts(stats, currency_stats).items()):
            col = index % chart_columns
            row = index // chart_columns
            left = chart_start_x + col * (chart_width + chart_gap_x)
            top = chart_start_y + row * (chart_row_height + chart_gap_y)
            slide.shapes.add_picture(BytesIO(png_bytes), left, top, width=chart_width)

    if report_type == "integral":
        _add_ngo_detail_slide(
            prs, blank_layout, stats, currency_stats, ngo_breakdown, sc, using_template
        )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_ngo_detail_slide(
    prs: Presentation,
    blank_layout,
    stats: DonationStatistics,
    currency_stats: CurrencyStatistics | None,
    ngo_breakdown: list[dict],
    sc,
    using_template: bool = False,
) -> None:
    slide = prs.slides.add_slide(blank_layout)

    title_text = f"Detail — Organisations funded ({_rotary_year_label(stats.selected_rotary_year)})"
    placeholder = _title_placeholder(slide) if using_template else None
    if placeholder is not None:
        placeholder.text_frame.text = title_text
    else:
        title_box = slide.shapes.add_textbox(
            sc(Inches(0.4)), sc(Inches(0.3)), sc(Inches(9)), sc(Inches(0.6))
        )
        title_tf = title_box.text_frame
        title_tf.text = title_text
        title_tf.paragraphs[0].font.size = Pt(20)
        title_tf.paragraphs[0].font.bold = True
        if not using_template:
            title_tf.paragraphs[0].font.color.rgb = RGBColor.from_string("17458F")

    if not ngo_breakdown or currency_stats is None:
        empty_box = slide.shapes.add_textbox(
            sc(Inches(0.4)), sc(Inches(1.1)), sc(Inches(9)), sc(Inches(0.5))
        )
        empty_box.text_frame.text = "No organisations funded this year."
        return

    headers = ["Organisation", "Total donated"]
    rows = [[row["name"], _format_currency(row["total"], currency_stats.currency)] for row in ngo_breakdown]

    table_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(
        table_rows, len(headers), sc(Inches(0.4)), sc(Inches(1.1)), sc(Inches(9)),
        sc(Inches(min(0.4 * table_rows, 5.5))),
    )
    table = table_shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for row_index, row_values in enumerate(rows, start=1):
        for col_index, value in enumerate(row_values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(10)
